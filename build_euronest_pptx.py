"""Euronest pitch deck — visual build.

16:9, vector icon system + stylised France map + native charts, soft
shadows, warm sand palette. Mixed media: photo slots on the cover,
solution and closing slides (drop files into photos/ and rebuild).

Copy written through the humanizer pass: no em dashes, no
"stands as / serves as", no rule-of-three by default, varied rhythm.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

from euronest_assets import icon, france_map, FRANCE_CITIES, FRANCE_VIEW

HERE = os.path.dirname(os.path.abspath(__file__))
PHOTO_DIR = os.path.join(HERE, "photos")

# ---- palette --------------------------------------------------------
NAVY   = RGBColor(0x10, 0x2A, 0x54)
NAVY2  = RGBColor(0x1B, 0x3A, 0x66)
ACCENT = RGBColor(0x3E, 0x63, 0xA6)
SKY    = RGBColor(0x6E, 0x8F, 0xC7)
SAND   = RGBColor(0xF5, 0xEF, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x22, 0x27, 0x33)
MUTED  = RGBColor(0x6B, 0x73, 0x80)
GOLD   = RGBColor(0xC9, 0x9A, 0x3C)
GOLDT  = RGBColor(0xF1, 0xE7, 0xCF)
ACCT   = RGBColor(0xE7, 0xED, 0xF6)
BORDER = RGBColor(0xE6, 0xDE, 0xCE)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF8)
GREEN  = RGBColor(0x4A, 0x7C, 0x59)
CORAL  = RGBColor(0xC4, 0x5A, 0x4E)

NAVY_HEX, ACCENT_HEX, WHITE_HEX = "102A54", "3E63A6", "FFFFFF"
GOLD_HEX, MUTED_HEX = "C99A3C", "6B7380"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ---- low-level helpers ---------------------------------------------
def soft_shadow(shape, blur=10, dist=4, color="9C8F76", alpha=42):
    """Attach a soft drop shadow to a shape."""
    sp = shape._element.spPr
    for tag in ("a:effectLst",):
        ex = sp.find(qn(tag))
        if ex is not None:
            sp.remove(ex)
    eff = sp.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Pt(blur)), "dist": str(Pt(dist)),
        "dir": "5400000", "rotWithShape": "0"})
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": color})
    a = clr.makeelement(qn("a:alpha"), {"val": str(alpha * 1000)})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    sp.append(eff)


def rect(slide, x, y, w, h, fill, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        soft_shadow(s)
    return s


def rrect(slide, x, y, w, h, fill, line=None, lw=0.75,
          radius=0.06, shadow=False):
    return rect(slide, x, y, w, h, fill, line, lw,
                MSO_SHAPE.ROUNDED_RECTANGLE, radius, shadow)


def oval(slide, x, y, w, h, fill, line=None, lw=1.0):
    return rect(slide, x, y, w, h, fill, line, lw, MSO_SHAPE.OVAL)


def text(slide, body, x, y, w, h, size=16, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         font="Calibri", spacing=1.15, space_after=0, shrink=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = body if isinstance(body, list) else body.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        runs = ln if isinstance(ln, list) else [(ln, {})]
        for seg, opt in runs:
            r = p.add_run()
            r.text = seg
            r.font.name = opt.get("font", font)
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.italic = opt.get("italic", italic)
            r.font.color.rgb = opt.get("color", color)
    return tb


def place_icon(slide, name, x, y, size, color_hex=NAVY_HEX, weight=1.65):
    p = icon(name, "#" + color_hex, size=260, weight=weight)
    slide.shapes.add_picture(p, x, y, size, size)


def icon_chip(slide, name, x, y, chip=Inches(0.82), fill=ACCT,
              icon_hex=ACCENT_HEX, radius=0.24, shadow=False):
    c = rrect(slide, x, y, chip, chip, fill, radius=radius, shadow=shadow)
    pad = chip * 0.24
    place_icon(slide, name, x + pad, y + pad, chip - 2 * pad, icon_hex)
    return c


def aspect_fill(pic, box_w, box_h, native_w, native_h):
    box_ar = box_w / box_h
    img_ar = native_w / native_h
    if img_ar > box_ar:
        crop = (1 - box_ar / img_ar) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ar / box_ar) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop


def place_photo(slide, key, x, y, w, h, label="Hero photo",
                spec="", rounded=False, overlay=False):
    """Place photos/<key>.jpg cropped to fill, or a styled placeholder."""
    path = None
    for ext in (".jpg", ".jpeg", ".png"):
        cand = os.path.join(PHOTO_DIR, key + ext)
        if os.path.exists(cand):
            path = cand
            break
    if path:
        from PIL import Image
        iw, ih = Image.open(path).size
        pic = slide.shapes.add_picture(path, x, y, w, h)
        aspect_fill(pic, w, h, iw, ih)
        if overlay:
            ov = rect(slide, x, y, w, h, NAVY)
            ov.fill.transparency = 0  # base
            _set_alpha(ov, 38)
        return pic
    # placeholder
    if rounded:
        ph = rrect(slide, x, y, w, h, NAVY2, radius=0.04)
    else:
        ph = rect(slide, x, y, w, h, NAVY2)
    ic = Inches(0.95)
    place_icon(slide, "camera", x + w / 2 - ic / 2,
               y + h / 2 - ic / 2 - Inches(0.35), ic, "8FA6CC")
    text(slide, label, x, y + h / 2 + Inches(0.35), w, Inches(0.4),
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if spec:
        text(slide, spec, x + Inches(0.3), y + h / 2 + Inches(0.78),
             w - Inches(0.6), Inches(0.8), size=10.5, color=SKY,
             align=PP_ALIGN.CENTER, spacing=1.25)
    return ph


def _set_alpha(shape, pct):
    """Set fill transparency (pct = how transparent, 0..100)."""
    sp = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if sp is not None:
        a = sp.makeelement(qn("a:alpha"), {"val": str((100 - pct) * 1000)})
        sp.append(a)


# ---- slide scaffolding ---------------------------------------------
PAGE = 0


def new_slide(dark=False):
    global PAGE
    PAGE += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY if dark else SAND)
    return s


def header(slide, title, eyebrow, num):
    rect(slide, 0, 0, SW, Inches(0.14), NAVY)
    text(slide, eyebrow.upper(), Inches(0.75), Inches(0.46),
         Inches(9), Inches(0.32), size=11.5, bold=True, color=ACCENT)
    text(slide, title, Inches(0.75), Inches(0.74),
         Inches(10.4), Inches(1.0), size=30, bold=True, color=NAVY)
    rect(slide, Inches(0.78), Inches(1.46), Inches(0.62),
         Inches(0.055), GOLD)
    # number chip
    ch = Inches(0.6)
    rrect(slide, SW - Inches(1.2), Inches(0.5), ch, ch, NAVY,
          radius=0.22)
    text(slide, f"{num:02d}", SW - Inches(1.2), Inches(0.5), ch, ch,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, light=True):
    col = MUTED if light else SKY
    text(slide, "EURONEST", Inches(0.75), Inches(7.04),
         Inches(3), Inches(0.3), size=9, bold=True, color=col)
    text(slide, "Student Housing & Lifestyle Subscription",
         Inches(2.0), Inches(7.04), Inches(6), Inches(0.3),
         size=9, color=col)
    text(slide, "April 2026", SW - Inches(2.0), Inches(7.04),
         Inches(1.25), Inches(0.3), size=9, color=col,
         align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — Cover
# =====================================================================
s = new_slide(dark=True)
# right photo panel
pw = Inches(5.0)
place_photo(s, "cover", SW - pw, 0, pw, SH,
            label="Cover photo",
            spec="Suggested: a bright Paris / Lyon street or a\n"
                 "furnished student apartment. Portrait or any\n"
                 "crop, min 1600x2000px.")
# accent seam
rect(s, SW - pw - Inches(0.08), 0, Inches(0.08), SH, GOLD)
# left content
rect(s, Inches(0.85), Inches(1.05), Inches(0.95), Inches(0.1), GOLD)
text(s, "EURONEST", Inches(0.8), Inches(1.25), Inches(7.2),
     Inches(1.5), size=72, bold=True, color=WHITE)
text(s, "Student housing and lifestyle,\nbuilt for the people who move.",
     Inches(0.85), Inches(2.95), Inches(7), Inches(1.3),
     size=21, color=SKY, italic=True, spacing=1.25)
# meta row
rect(s, Inches(0.85), Inches(4.35), Inches(6.6), Inches(0.018), NAVY2)
text(s, "BUSINESS PITCH", Inches(0.85), Inches(4.55),
     Inches(4), Inches(0.35), size=12, bold=True, color=GOLD)
text(s, "ISTEC, Paris  ·  April 2026", Inches(0.85), Inches(4.85),
     Inches(6), Inches(0.4), size=14, color=WHITE)
text(s,
     "Mansoor Ali   ·   Fidha Fathima Panankavil   ·   Leo Joy\n"
     "Nithin Stanley George   ·   Minhaj Eswaramangalam\n"
     "Gopinath Dasan   ·   Nithin Pallampetti Velayudhan",
     Inches(0.85), Inches(5.5), Inches(7.2), Inches(1.3),
     size=12.5, color=SKY, spacing=1.45)


# =====================================================================
# SLIDE 2 — Agenda
# =====================================================================
s = new_slide()
header(s, "What we will cover", "Agenda", PAGE)
agenda = [
    ("search", "The problem", "Why housing in France fails students"),
    ("shield_check", "Our solution", "What Euronest actually does"),
    ("globe", "The market", "Where we start and who we serve"),
    ("wallet", "Business model", "Tiers, revenue and unit economics"),
    ("trend", "Financials & ask", "Three-year plan and the EUR 180k raise"),
    ("rocket", "Team & roadmap", "Who builds this, and what comes next"),
]
cw, chh = Inches(5.9), Inches(1.32)
gx, gy = Inches(0.75), Inches(2.0)
for i, (ic, t, sub) in enumerate(agenda):
    r, c = i // 2, i % 2
    x = gx + c * (cw + Inches(0.35))
    y = gy + r * (chh + Inches(0.28))
    card = rrect(s, x, y, cw, chh, WHITE, line=BORDER,
                 radius=0.1, shadow=True)
    icon_chip(s, ic, x + Inches(0.3), y + Inches(0.28),
              chip=Inches(0.76))
    text(s, f"{i+1:02d}", x + cw - Inches(1.05), y + Inches(0.2),
         Inches(0.9), Inches(0.5), size=24, bold=True, color=ACCT,
         align=PP_ALIGN.RIGHT)
    text(s, t, x + Inches(1.32), y + Inches(0.28),
         cw - Inches(1.5), Inches(0.45), size=17, bold=True, color=NAVY)
    text(s, sub, x + Inches(1.32), y + Inches(0.72),
         cw - Inches(1.5), Inches(0.45), size=12, color=MUTED)
footer(s)


# =====================================================================
# SLIDE 3 — Problem
# =====================================================================
s = new_slide()
header(s, "Housing in France is broken for students",
       "01 · The problem", PAGE)
probs = [
    ("alert", "1 in 3", "students we interviewed had paid money to a "
     "fake or misleading listing."),
    ("doc", "20 pages", "of dense French legal lease, usually signed "
     "without fully understanding it."),
    ("euro", "3x rent", "the income a French guarantor must show. Most "
     "non-EU students cannot find one."),
    ("calendar", "6 months", "rent demanded upfront when the guarantor "
     "wall blocks a student."),
]
cw = Inches(2.83)
gx, gy, chh = Inches(0.78), Inches(2.0), Inches(3.5)
for i, (ic, big, body) in enumerate(probs):
    x = gx + i * (cw + Inches(0.17))
    card = rrect(s, x, gy, cw, chh, WHITE, line=BORDER,
                 radius=0.08, shadow=True)
    rrect(s, x, gy, cw, Inches(0.12), CORAL, radius=0.0)
    icon_chip(s, ic, x + Inches(0.32), gy + Inches(0.34),
              chip=Inches(0.84), fill=RGBColor(0xF6, 0xE3, 0xE1),
              icon_hex="C45A4E")
    text(s, big, x + Inches(0.32), gy + Inches(1.42),
         cw - Inches(0.6), Inches(0.7), size=30, bold=True, color=NAVY)
    text(s, body, x + Inches(0.32), gy + Inches(2.12),
         cw - Inches(0.6), Inches(1.2), size=12.5, color=INK,
         spacing=1.3)
# source strip
rrect(s, Inches(0.75), Inches(5.78), Inches(11.83), Inches(0.62),
      ACCT, radius=0.18)
text(s, "Source: 42 interviews with international students at ISTEC "
        "and three other Paris schools, January to March 2026.",
     Inches(1.1), Inches(5.78), Inches(11.2), Inches(0.62),
     size=11.5, italic=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# =====================================================================
# SLIDE 4 — Why now
# =====================================================================
s = new_slide()
header(s, "The timing is on our side", "01 · Why now", PAGE)
facts = [
    ("trend", "430,000", "international students in France, 2024-25.",
     "Government target: 500,000 by 2027."),
    ("building", "~175,000", "publicly owned CROUS rooms.",
     "Long waitlists, allocation favours French and EU students."),
    ("scatter", "under 15%", "of private student stock is professionally "
     "managed.", "The rest is scattered across small landlords."),
]
gy = Inches(2.05)
for i, (ic, big, line, sub) in enumerate(facts):
    y = gy + i * Inches(1.34)
    icon_chip(s, ic, Inches(0.85), y, chip=Inches(0.98), shadow=True)
    text(s, big, Inches(2.05), y - Inches(0.04), Inches(2.5),
         Inches(0.7), size=33, bold=True, color=ACCENT)
    text(s, line, Inches(4.55), y + Inches(0.03), Inches(4.2),
         Inches(0.8), size=14.5, bold=True, color=NAVY, spacing=1.15)
    text(s, sub, Inches(4.55), y + Inches(0.6), Inches(4.3),
         Inches(0.6), size=11.5, color=MUTED)
# right insight panel
px = Inches(9.15)
panel = rrect(s, px, gy - Inches(0.05), Inches(3.45), Inches(4.4),
              NAVY, radius=0.07, shadow=True)
place_icon(s, "compass", px + Inches(0.32), gy + Inches(0.25),
           Inches(0.6), GOLD_HEX)
text(s, "Why this window", px + Inches(0.32), gy + Inches(1.0),
     Inches(2.9), Inches(0.4), size=13, bold=True, color=GOLD)
text(s, "The 2026 intake grew up booking flights, food and rides "
        "from an app.\n\nThey expect housing to work the same way.\n\n"
        "Right now, it does not.",
     px + Inches(0.32), gy + Inches(1.45), Inches(2.85), Inches(2.8),
     size=13.5, color=WHITE, spacing=1.4)
footer(s)


# =====================================================================
# SLIDE 5 — Solution
# =====================================================================
s = new_slide()
header(s, "Euronest, in one frame", "02 · The solution", PAGE)
pillars = [
    ("shield_check", "Verified apartment",
     "Every property is visited and checked against a 38-point list. "
     "We refuse roughly a third of what landlords offer."),
    ("wallet", "One bundled price",
     "Rent, utilities, internet and lifestyle services on a single "
     "invoice. Euronest is the guarantor on standard and premium."),
    ("headset", "One named coordinator",
     "Each subscriber gets one multilingual contact, reachable on "
     "WhatsApp from move-in to move-out."),
]
cw = Inches(2.62)
gx, gy, chh = Inches(0.75), Inches(2.0), Inches(4.35)
for i, (ic, t, body) in enumerate(pillars):
    x = gx + i * (cw + Inches(0.2))
    rrect(s, x, gy, cw, chh, WHITE, line=BORDER, radius=0.07,
          shadow=True)
    icon_chip(s, ic, x + Inches(0.32), gy + Inches(0.34),
              chip=Inches(0.92), shadow=False)
    text(s, f"{i+1}", x + cw - Inches(0.85), gy + Inches(0.3),
         Inches(0.6), Inches(0.6), size=22, bold=True, color=ACCT)
    text(s, t, x + Inches(0.32), gy + Inches(1.5),
         cw - Inches(0.6), Inches(0.8), size=16.5, bold=True,
         color=NAVY, spacing=1.1)
    text(s, body, x + Inches(0.32), gy + Inches(2.35),
         cw - Inches(0.6), Inches(1.8), size=12, color=INK,
         spacing=1.32)
# right photo panel
px = Inches(9.0)
place_photo(s, "solution", px, gy, Inches(3.6), chh,
            label="Lifestyle photo", rounded=True,
            spec="Suggested: a student at home, cooking,\n"
                 "or with friends. Portrait crop,\nmin 1200x1500px.")
footer(s)


# =====================================================================
# SLIDE 6 — How it works
# =====================================================================
s = new_slide()
header(s, "How a student becomes a subscriber",
       "02 · User journey", PAGE)
steps = [
    ("search", "Discover", "A TikTok, a friend, or a YouTube city "
     "guide. They land on the site and watch a 90-second explainer."),
    ("cursor", "Book", "They pick a city, a tier and a move-in date, "
     "sign a bilingual lease, and pay into escrow."),
    ("key", "Arrive", "The coordinator meets them at the door for an "
     "inventory walkthrough. Wi-Fi is already on."),
    ("home", "Live", "WhatsApp stays open. Maintenance SLAs are "
     "published. The quiet-hours code is signed."),
    ("refresh", "Renew", "Two months out: renew, swap city, or get "
     "matched with a roommate for next year."),
]
n = len(steps)
gx, gy = Inches(0.75), Inches(2.35)
cw = Inches(2.30)
gap = Inches(0.07)
# connector line
rect(s, gx + cw / 2, gy + Inches(0.46), (cw + gap) * (n - 1),
     Inches(0.05), SKY)
for i, (ic, t, body) in enumerate(steps):
    x = gx + i * (cw + gap)
    cx = x + cw / 2
    disc = Inches(0.96)
    oval(s, cx - disc / 2, gy, disc, disc, NAVY)
    place_icon(s, ic, cx - disc / 2 + Inches(0.24),
               gy + Inches(0.24), Inches(0.48), WHITE_HEX)
    # step number badge
    bd = Inches(0.36)
    oval(s, cx + Inches(0.18), gy - Inches(0.06), bd, bd, GOLD)
    text(s, str(i + 1), cx + Inches(0.18), gy - Inches(0.06), bd, bd,
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    card = rrect(s, x, gy + Inches(1.3), cw, Inches(2.75), WHITE,
                 line=BORDER, radius=0.09, shadow=True)
    text(s, t, x, gy + Inches(1.5), cw, Inches(0.45), size=16,
         bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s, cx - Inches(0.2), gy + Inches(1.95), Inches(0.4),
         Inches(0.04), GOLD)
    text(s, body, x + Inches(0.22), gy + Inches(2.12),
         cw - Inches(0.44), Inches(1.8), size=11.3, color=INK,
         align=PP_ALIGN.CENTER, spacing=1.3)
footer(s)


# =====================================================================
# SLIDE 7 — Tiers
# =====================================================================
s = new_slide()
header(s, "Three tiers, one monthly invoice", "02 · The offer", PAGE)
tiers = [
    ("Basic", "690", NAVY, False,
     ["Furnished apartment", "Fibre internet", "Utilities included",
      "Contents insurance"],
     "Short stays and budget-minded Erasmus students."),
    ("Standard", "890", ACCENT, True,
     ["Everything in Basic", "Gym and fitness pass",
      "Netflix or Amazon Prime", "Euronest acts as guarantor"],
     "Master's students staying a full academic year."),
    ("Premium", "1,190", GOLD, False,
     ["Everything in Standard", "Meal plan, 15 meals a week",
      "Visa and paperwork support", "Annual return-flight credit"],
     "Students whose families want the full setup."),
]
cw = Inches(3.85)
gx, gy = Inches(0.78), Inches(1.95)
for i, (name, price, color, featured, feats, who) in enumerate(tiers):
    x = gx + i * (cw + Inches(0.18))
    ch = Inches(4.7) if featured else Inches(4.45)
    yy = gy - (Inches(0.13) if featured else 0)
    card = rrect(s, x, yy, cw, ch, WHITE, line=BORDER,
                 radius=0.06, shadow=True)
    head = rrect(s, x, yy, cw, Inches(1.18), color, radius=0.06)
    rect(s, x, yy + Inches(0.6), cw, Inches(0.58), color)
    if featured:
        tag = rrect(s, x + cw - Inches(1.7), yy + Inches(0.2),
                    Inches(1.5), Inches(0.36), GOLD, radius=0.5)
        text(s, "MOST POPULAR", x + cw - Inches(1.7),
             yy + Inches(0.2), Inches(1.5), Inches(0.36), size=8.5,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text(s, name, x + Inches(0.35), yy + Inches(0.16),
         Inches(2.5), Inches(0.45), size=18, bold=True, color=WHITE)
    text(s, [[("EUR ", {"size": 13, "color": WHITE}),
              (price, {"size": 30, "bold": True, "color": WHITE}),
              (" / mo", {"size": 13, "color": WHITE})]],
         x + Inches(0.35), yy + Inches(0.58), cw - Inches(0.6),
         Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    fy = yy + Inches(1.45)
    for f in feats:
        place_icon(s, "check", x + Inches(0.35), fy + Inches(0.02),
                   Inches(0.26), ACCENT_HEX)
        text(s, f, x + Inches(0.75), fy, cw - Inches(1.0),
             Inches(0.4), size=12, color=INK)
        fy += Inches(0.5)
    rect(s, x + Inches(0.35), fy + Inches(0.05), cw - Inches(0.7),
         Inches(0.02), BORDER)
    text(s, who, x + Inches(0.35), fy + Inches(0.18),
         cw - Inches(0.7), Inches(0.8), size=10.8, italic=True,
         color=MUTED, spacing=1.25)
text(s, "Paris pricing. Other cities run 15-25% lower. Six-month "
        "prepay carries a 5% discount.",
     Inches(0.78), Inches(6.82), Inches(11.8), Inches(0.35),
     size=10.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
footer(s)


# =====================================================================
# SLIDE 8 — Market (France map)
# =====================================================================
s = new_slide()
header(s, "Five cities, about 60% of the market", "03 · The market",
       PAGE)
# France map
fpath, (vw, vh) = france_map()
map_h = Inches(4.55)
map_w = map_h * vw / vh
mx, my = Inches(0.85), Inches(1.95)
s.shapes.add_picture(fpath, mx, my, map_w, map_h)
cities = {
    "Paris": ("Year 1", GOLD), "Lyon": ("Year 1", GOLD),
    "Lille": ("Year 2", ACCENT), "Toulouse": ("Year 2", ACCENT),
    "Marseille": ("Year 3", SKY),
}
for name, (cx, cy) in FRANCE_CITIES.items():
    px = mx + Emu(int(map_w * cx / vw))
    py = my + Emu(int(map_h * cy / vh))
    wave, col = cities[name]
    d = Inches(0.26)
    oval(s, px - d / 2, py - d / 2, d, d, col, line=WHITE, lw=1.5)
    lab_x = px + Inches(0.18)
    text(s, name, lab_x, py - Inches(0.17), Inches(1.6),
         Inches(0.34), size=12, bold=True, color=NAVY)
# right: city list
lx, ly = Inches(6.6), Inches(2.0)
rows = [
    ("Paris", "115,000", "Year 1", GOLD),
    ("Lyon", "32,000", "Year 1", GOLD),
    ("Lille", "21,000", "Year 2", ACCENT),
    ("Toulouse", "19,000", "Year 2", ACCENT),
    ("Marseille", "17,000", "Year 3", SKY),
]
text(s, [[("City", {}), ]], lx + Inches(0.25), ly, Inches(2),
     Inches(0.3), size=11, bold=True, color=MUTED)
text(s, "INTL. STUDENTS", lx + Inches(2.5), ly, Inches(2),
     Inches(0.3), size=11, bold=True, color=MUTED)
text(s, "ENTRY", lx + Inches(4.5), ly, Inches(1.3), Inches(0.3),
     size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
ry = ly + Inches(0.4)
for name, num, wave, col in rows:
    rrect(s, lx, ry, Inches(5.9), Inches(0.62), WHITE, line=BORDER,
          radius=0.16)
    rect(s, lx, ry, Inches(0.1), Inches(0.62), col)
    text(s, name, lx + Inches(0.3), ry, Inches(2.2), Inches(0.62),
         size=13.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, num, lx + Inches(2.5), ry, Inches(2), Inches(0.62),
         size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    tag = rrect(s, lx + Inches(4.55), ry + Inches(0.13),
                Inches(1.15), Inches(0.36), col, radius=0.5)
    text(s, wave, lx + Inches(4.55), ry + Inches(0.13),
         Inches(1.15), Inches(0.36), size=9.5, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ry += Inches(0.72)
text(s, "Paris and Lyon at launch give us critical mass without "
        "spreading the team too thin. Lille and Toulouse open the "
        "Erasmus corridor; Marseille closes the southern triangle.",
     lx, ry + Inches(0.05), Inches(5.9), Inches(1.0), size=11.5,
     color=MUTED, spacing=1.3, italic=True)
footer(s)


# =====================================================================
# SLIDE 9 — Competition
# =====================================================================
s = new_slide()
header(s, "Nobody else bundles all three", "03 · Competition", PAGE)
# matrix
mx, my = Inches(0.85), Inches(2.0)
mw, mh = Inches(7.1), Inches(4.45)
rrect(s, mx, my, mw, mh, WHITE, line=BORDER, radius=0.03, shadow=True)
# quadrant tint (top-right = our zone)
rect(s, mx + mw / 2, my, mw / 2, mh / 2, ACCT)
# axes
rect(s, mx + mw / 2 - Inches(0.01), my, Inches(0.025), mh, BORDER)
rect(s, mx, my + mh / 2 - Inches(0.01), mw, Inches(0.025), BORDER)
text(s, "MORE SERVICE  ▲", mx + Inches(0.15), my + Inches(0.12),
     Inches(3), Inches(0.3), size=9.5, bold=True, color=MUTED)
text(s, "MORE INVENTORY  ▶", mx + mw - Inches(2.5),
     my + mh - Inches(0.38), Inches(2.4), Inches(0.3), size=9.5,
     bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

def plot(xp, yp, label, col=ACCENT, big=False):
    cx = mx + Emu(int(mw * xp))
    cy = my + Emu(int(mh * (1 - yp)))
    d = Inches(0.34) if big else Inches(0.2)
    oval(s, cx - d / 2, cy - d / 2, d, d, col,
         line=WHITE, lw=1.5)
    text(s, label, cx + d / 2 + Inches(0.06), cy - Inches(0.15),
         Inches(2.4), Inches(0.32), size=11,
         bold=big, color=NAVY if big else INK)

plot(0.16, 0.18, "Airbnb", SKY)
plot(0.34, 0.30, "Studapart", SKY)
plot(0.40, 0.46, "HousingAnywhere", SKY)
plot(0.50, 0.56, "Student.com", SKY)
plot(0.58, 0.50, "Uniplaces", SKY)
plot(0.12, 0.10, "CROUS", SKY)
plot(0.22, 0.52, "Local agencies", SKY)
plot(0.78, 0.84, "Euronest", GOLD, big=True)
# right takeaway
px = Inches(8.35)
text(s, "Our position", px, Inches(2.1), Inches(4), Inches(0.4),
     size=17, bold=True, color=NAVY)
rect(s, px, Inches(2.5), Inches(0.5), Inches(0.05), GOLD)
for i, (t, b) in enumerate([
    ("Curated inventory", "We hand-pick and verify, instead of "
     "listing everything."),
    ("In-house guarantor", "Backed by our own balance sheet, not a "
     "third party."),
    ("Service layer", "A bundle that only pays off if you specialise "
     "in students."),
]):
    yy = Inches(2.8) + i * Inches(1.05)
    place_icon(s, "check", px, yy + Inches(0.02), Inches(0.3),
               GOLD_HEX)
    text(s, t, px + Inches(0.45), yy - Inches(0.04), Inches(3.8),
         Inches(0.35), size=13.5, bold=True, color=NAVY)
    text(s, b, px + Inches(0.45), yy + Inches(0.3), Inches(3.85),
         Inches(0.7), size=11, color=MUTED, spacing=1.25)
text(s, "Pieces of our offer exist everywhere. Assembling all "
        "three is the moat.",
     px, Inches(6.0), Inches(4.2), Inches(0.7), size=11.5,
     italic=True, color=ACCENT, spacing=1.3)
footer(s)


# =====================================================================
# SLIDE 10 — Business model
# =====================================================================
s = new_slide()
header(s, "How the money works", "04 · Business model", PAGE)
# donut chart — revenue mix
cd = CategoryChartData()
cd.categories = ["Subscriptions", "Landlord commission",
                 "Service partners", "Add-ons"]
cd.add_series("Revenue mix", (90, 5, 3, 2))
gframe = s.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(0.75), Inches(1.95),
    Inches(4.5), Inches(4.4), cd)
chart = gframe.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
plot = chart.plots[0]
plot.donut_hole_size = 58
dcols = [NAVY, ACCENT, SKY, GOLD]
for i, pt in enumerate(plot.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = dcols[i]
    pt.format.line.color.rgb = SAND
    pt.format.line.width = Pt(2)
text(s, "Revenue mix", Inches(0.75), Inches(1.95), Inches(4.5),
     Inches(0.4), size=13, bold=True, color=NAVY,
     align=PP_ALIGN.CENTER)
# revenue stream rows
sx = Inches(5.5)
streams = [
    ("euro", "Subscriptions", "Monthly or six-month prepay across "
     "the three tiers.", NAVY),
    ("handshake", "Landlord commission", "6% on the first month for "
     "listings we route but do not head-lease.", ACCENT),
    ("dumbbell", "Service partners", "A fixed margin from gyms, ISPs "
     "and meal providers.", SKY),
    ("plus", "Add-ons", "Airport pickup, SIM cards, language "
     "coaching.", GOLD),
]
sy = Inches(2.0)
for ic, t, b, col in streams:
    icon_chip(s, ic, sx, sy, chip=Inches(0.66),
              fill=ACCT, icon_hex=ACCENT_HEX)
    text(s, t, sx + Inches(0.85), sy - Inches(0.02), Inches(3.2),
         Inches(0.35), size=13.5, bold=True, color=NAVY)
    text(s, b, sx + Inches(0.85), sy + Inches(0.3), Inches(3.3),
         Inches(0.6), size=10.8, color=MUTED, spacing=1.2)
    sy += Inches(0.92)
# unit economics card
ux = Inches(9.35)
rrect(s, ux, Inches(2.0), Inches(3.25), Inches(4.3), NAVY,
      radius=0.06, shadow=True)
text(s, "Unit economics", ux + Inches(0.3), Inches(2.2),
     Inches(2.7), Inches(0.4), size=13, bold=True, color=GOLD)
text(s, "Year 1, blended", ux + Inches(0.3), Inches(2.5),
     Inches(2.7), Inches(0.3), size=10, color=SKY)
ue = [("Avg. subscription", "EUR 870"), ("Cost of services", "EUR 390"),
      ("Contribution", "EUR 480"), ("Fixed monthly opex", "EUR 18.5k"),
      ("Break-even", "39 subs")]
uy = Inches(2.95)
for i, (k, v) in enumerate(ue):
    text(s, k, ux + Inches(0.3), uy, Inches(1.9), Inches(0.4),
         size=11.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, v, ux + Inches(2.05), uy, Inches(0.95), Inches(0.4),
         size=11.5, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE)
    if i < len(ue) - 1:
        rect(s, ux + Inches(0.3), uy + Inches(0.5), Inches(2.65),
             Pt(0.75), NAVY2)
    uy += Inches(0.62)
footer(s)


# =====================================================================
# SLIDE 11 — Service standards
# =====================================================================
s = new_slide()
header(s, "What we promise the subscriber",
       "04 · Service standards", PAGE)
sla = [
    ("alert", "Emergency", "Within 24 hours",
     "No power, water leak, lock failure.", CORAL,
     RGBColor(0xF6, 0xE3, 0xE1), "C45A4E"),
    ("clock", "Urgent", "Within 5 working days",
     "Broken appliance, heater, prolonged outage.", GOLD,
     GOLDT, "B5862F"),
    ("wrench", "Non-urgent", "Within 28 days",
     "Single fitting, sticking door or window.", GREEN,
     RGBColor(0xE2, 0xED, 0xE5), "3E6B4C"),
]
cw = Inches(3.85)
gx, gy = Inches(0.78), Inches(1.95)
for i, (ic, name, when, ex, bar, tint, ihex) in enumerate(sla):
    x = gx + i * (cw + Inches(0.18))
    rrect(s, x, gy, cw, Inches(2.5), WHITE, line=BORDER,
          radius=0.08, shadow=True)
    rrect(s, x, gy, Inches(0.14), Inches(2.5), bar, radius=0.0)
    icon_chip(s, ic, x + Inches(0.34), gy + Inches(0.3),
              chip=Inches(0.78), fill=tint, icon_hex=ihex)
    text(s, name, x + Inches(1.3), gy + Inches(0.32),
         cw - Inches(1.5), Inches(0.4), size=16, bold=True,
         color=NAVY)
    text(s, when, x + Inches(1.3), gy + Inches(0.68),
         cw - Inches(1.5), Inches(0.4), size=12, bold=True,
         color=bar)
    text(s, ex, x + Inches(0.34), gy + Inches(1.32),
         cw - Inches(0.7), Inches(1.0), size=11.5, color=INK,
         spacing=1.3)
# bottom band
by = Inches(4.75)
rrect(s, Inches(0.78), by, Inches(11.77), Inches(1.85), NAVY,
      radius=0.05, shadow=True)
text(s, "Also written into every lease", Inches(1.1),
     by + Inches(0.22), Inches(10), Inches(0.4), size=13,
     bold=True, color=GOLD)
commits = [
    ("person", "One named, multilingual coordinator. No call centre."),
    ("moon", "Quiet hours 22:00-07:00, four-step escalation."),
    ("fire", "Smoke alarm, blanket and extinguisher in every unit."),
    ("doc", "Bilingual lease, deposit held in escrow."),
]
cw2 = Inches(2.85)
for i, (ic, t) in enumerate(commits):
    x = Inches(1.1) + i * cw2
    place_icon(s, ic, x, by + Inches(0.72), Inches(0.42), GOLD_HEX)
    text(s, t, x + Inches(0.56), by + Inches(0.68),
         cw2 - Inches(0.7), Inches(0.95), size=10.8, color=WHITE,
         spacing=1.25)
footer(s)


# =====================================================================
# SLIDE 12 — Financials
# =====================================================================
s = new_slide()
header(s, "Three-year projection", "05 · Financials", PAGE)
cd = CategoryChartData()
cd.categories = ["Year 1", "Year 2", "Year 3"]
cd.add_series("Revenue", (540, 1820, 3360))
cd.add_series("Gross margin", (155, 610, 1180))
cd.add_series("EBITDA", (-35, 215, 540))
gframe = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(2.0),
    Inches(7.5), Inches(4.4), cd)
chart = gframe.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
for i, ser in enumerate(chart.plots[0].series):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = [NAVY, ACCENT, GOLD][i]
    ser.format.line.fill.background()
try:
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(11)
except Exception:
    pass
text(s, "EUR thousands", Inches(0.75), Inches(1.95), Inches(3),
     Inches(0.3), size=9.5, italic=True, color=MUTED)
# narrative panel
px = Inches(8.55)
rrect(s, px, Inches(2.0), Inches(4.0), Inches(4.4), NAVY,
      radius=0.06, shadow=True)
place_icon(s, "trend", px + Inches(0.3), Inches(2.2), Inches(0.55),
           GOLD_HEX)
text(s, "Reading the numbers", px + Inches(0.3), Inches(2.9),
     Inches(3.4), Inches(0.4), size=13, bold=True, color=GOLD)
for i, (t, b) in enumerate([
    ("Year 1 is a build year", "The intake discount and a full team "
     "from month six put EBITDA at -EUR 35k."),
    ("Positive in month 13-15", "Depending on how June to August "
     "occupancy lands."),
    ("Year 3", "EBITDA of EUR 540k on EUR 3.36m revenue, a 16% "
     "margin."),
]):
    yy = Inches(3.35) + i * Inches(1.0)
    text(s, t, px + Inches(0.3), yy, Inches(3.45), Inches(0.35),
         size=12, bold=True, color=WHITE)
    text(s, b, px + Inches(0.3), yy + Inches(0.32), Inches(3.45),
         Inches(0.65), size=10.5, color=SKY, spacing=1.25)
footer(s)


# =====================================================================
# SLIDE 13 — The ask
# =====================================================================
s = new_slide()
header(s, "EUR 180,000 to launch", "05 · The ask", PAGE)
# left hero block
bx, by = Inches(0.78), Inches(2.0)
rrect(s, bx, by, Inches(5.3), Inches(4.35), NAVY, radius=0.06,
      shadow=True)
place_icon(s, "euro", bx + Inches(0.4), by + Inches(0.38),
           Inches(0.7), GOLD_HEX)
text(s, [[("EUR ", {"size": 26, "color": WHITE, "bold": True}),
          ("180k", {"size": 58, "color": WHITE, "bold": True})]],
     bx + Inches(0.4), by + Inches(1.15), Inches(4.6), Inches(1.1),
     anchor=MSO_ANCHOR.MIDDLE)
text(s, "initial round", bx + Inches(0.45), by + Inches(2.2),
     Inches(4.5), Inches(0.4), size=15, color=GOLD)
rect(s, bx + Inches(0.45), by + Inches(2.7), Inches(4.4),
     Pt(1), NAVY2)
for i, (k, v) in enumerate([
    ("Founders", "EUR 40k"),
    ("Bpifrance student loan", "EUR 60k"),
    ("Angel round (12% equity)", "EUR 80k")]):
    yy = by + Inches(2.95) + i * Inches(0.44)
    text(s, k, bx + Inches(0.45), yy, Inches(3.1), Inches(0.4),
         size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, v, bx + Inches(3.4), yy, Inches(1.45), Inches(0.4),
         size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE)
# right: use of funds
ux = Inches(6.5)
text(s, "Where it goes", ux, Inches(2.0), Inches(6), Inches(0.45),
     size=17, bold=True, color=NAVY)
rect(s, ux, Inches(2.45), Inches(0.5), Inches(0.05), GOLD)
uses = [
    ("Furniture, deposits, setup", 55, "EUR 55k"),
    ("Six months of fixed opex", 45, "EUR 45k"),
    ("Platform and mobile app", 35, "EUR 35k"),
    ("September 2026 marketing", 25, "EUR 25k"),
    ("Working capital reserve", 20, "EUR 20k"),
]
uy = Inches(2.75)
maxv = 55
barmax = Inches(5.7)
for label, val, disp in uses:
    text(s, label, ux, uy, Inches(4.2), Inches(0.35), size=12,
         color=INK)
    text(s, disp, ux + barmax - Inches(1.2), uy, Inches(1.2),
         Inches(0.35), size=12, bold=True, color=NAVY,
         align=PP_ALIGN.RIGHT)
    rrect(s, ux, uy + Inches(0.34), barmax, Inches(0.2), LIGHT,
          radius=0.5)
    rrect(s, ux, uy + Inches(0.34), Emu(int(barmax * val / maxv)),
          Inches(0.2), ACCENT, radius=0.5)
    uy += Inches(0.72)
text(s, "Half the raise goes into deposits and furniture, which we "
        "recover as leases turn over.",
     ux, uy + Inches(0.05), Inches(5.7), Inches(0.6), size=11,
     italic=True, color=MUTED, spacing=1.3)
footer(s)


# =====================================================================
# SLIDE 14 — Team
# =====================================================================
s = new_slide()
header(s, "The founding team", "06 · Team", PAGE)
team = [
    ("Mansoor Ali", "CEO, partnerships",
     "Hospitality operations; ran a 40-room hotel.", NAVY),
    ("Fidha Fathima Panankavil", "COO, apartments",
     "Real estate analyst; inventory and setup.", ACCENT),
    ("Leo Joy", "CFO",
     "Audit and corporate finance; the model and the raise.", SKY),
    ("Nithin Stanley George", "CTO, platform",
     "Software engineer; website, app and payments.", NAVY),
    ("Minhaj Eswaramangalam", "Head of marketing",
     "Digital marketing; paid social and referrals.", ACCENT),
    ("Gopinath Dasan", "Head of support",
     "Multilingual support; leads the coordinators.", SKY),
    ("Nithin Pallampetti Velayudhan", "Head of legal",
     "Law graduate; leases, GDPR, compliance.", NAVY),
]


def initials(name):
    parts = [p for p in name.split() if p]
    return (parts[0][0] + parts[-1][0]).upper()


cw, chh = Inches(2.83), Inches(2.42)
gx, gy = Inches(0.78), Inches(1.85)
gapx, gapy = Inches(0.13), Inches(0.14)
for i, (name, role, bg, col) in enumerate(team):
    r, c = i // 4, i % 4
    x = gx + c * (cw + gapx)
    y = gy + r * (chh + gapy)
    rrect(s, x, y, cw, chh, WHITE, line=BORDER, radius=0.07,
          shadow=True)
    av = Inches(0.58)
    oval(s, x + Inches(0.24), y + Inches(0.22), av, av, col)
    text(s, initials(name), x + Inches(0.24), y + Inches(0.22),
         av, av, size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, role.upper(), x + Inches(0.95), y + Inches(0.22),
         cw - Inches(1.2), Inches(0.58), size=8.5, bold=True,
         color=ACCENT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
    text(s, name, x + Inches(0.24), y + Inches(0.92),
         cw - Inches(0.48), Inches(0.58), size=12.5, bold=True,
         color=NAVY, spacing=1.05)
    rect(s, x + Inches(0.24), y + Inches(1.58), cw - Inches(0.48),
         Pt(0.75), BORDER)
    text(s, bg, x + Inches(0.24), y + Inches(1.7),
         cw - Inches(0.48), Inches(0.62), size=10, color=MUTED,
         spacing=1.22)
# 8th cell — summary
x = gx + 3 * (cw + gapx)
y = gy + 1 * (chh + gapy)
rrect(s, x, y, cw, chh, NAVY, radius=0.07, shadow=True)
place_icon(s, "users", x + Inches(0.24), y + Inches(0.26),
           Inches(0.52), GOLD_HEX)
text(s, "Seven founders.", x + Inches(0.24), y + Inches(1.0),
     cw - Inches(0.48), Inches(0.4), size=14, bold=True,
     color=WHITE)
text(s, "Every one of us has been the customer we are now "
        "building for.",
     x + Inches(0.24), y + Inches(1.38), cw - Inches(0.48),
     Inches(0.9), size=11, color=SKY, spacing=1.3)
footer(s)


# =====================================================================
# SLIDE 15 — Risks
# =====================================================================
s = new_slide()
header(s, "Risks we are watching", "06 · Risk", PAGE)
risks = [
    ("New brand, low trust", "HIGH", CORAL,
     "Money-back first month, verified reviews, 10% first-cohort "
     "discount for September 2026."),
    ("Cash gap before break-even", "MEDIUM", GOLD,
     "Six-month prepay pulls cash forward; a EUR 60k bridge is "
     "pre-agreed with the angel."),
    ("Regulation on furnished lets", "LOW / HIGH", ACCENT,
     "Quarterly legal review and FNAIM membership for early "
     "warning."),
    ("Apartment damage or non-payment", "MEDIUM", GOLD,
     "Two-month deposit in escrow; liability cover bundled into "
     "the standard tier."),
    ("A key founder leaves early", "MEDIUM", GOLD,
     "Four-year vesting with a one-year cliff; two-deep cover on "
     "every function."),
]
gy = Inches(2.0)
for i, (t, lvl, col, mit) in enumerate(risks):
    y = gy + i * Inches(0.93)
    rrect(s, Inches(0.78), y, Inches(11.77), Inches(0.8), WHITE,
          line=BORDER, radius=0.14, shadow=True)
    rect(s, Inches(0.78), y, Inches(0.12), Inches(0.8), col)
    place_icon(s, "alert", Inches(1.08), y + Inches(0.2),
               Inches(0.4), MUTED_HEX)
    text(s, t, Inches(1.7), y, Inches(3.5), Inches(0.8), size=13,
         bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    tag = rrect(s, Inches(5.3), y + Inches(0.22), Inches(1.5),
                Inches(0.36), col, radius=0.5)
    text(s, lvl, Inches(5.3), y + Inches(0.22), Inches(1.5),
         Inches(0.36), size=9, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, mit, Inches(7.05), y, Inches(5.3), Inches(0.8),
         size=10.8, color=INK, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)
text(s, "Our own stop rule: under 120 subscribers and CAC above "
        "EUR 280 at month 18, we wind down.",
     Inches(0.78), Inches(6.75), Inches(11.8), Inches(0.35),
     size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
footer(s)


# =====================================================================
# SLIDE 16 — Roadmap
# =====================================================================
s = new_slide()
header(s, "The road from here", "06 · Roadmap", PAGE)
miles = [
    ("Q2 2026", "flag", "Incorporate, close the raise, ship the "
     "booking website."),
    ("Q3 2026", "key", "First 25 apartments live. September intake: "
     "20 paid subscribers."),
    ("Q1 2027", "users", "55 active subscribers. Mobile app on iOS "
     "and Android."),
    ("Q3 2027", "pin", "Lille and Toulouse open. AI roommate "
     "matching goes live."),
    ("Q4 2027", "trend", "180 subscribers. First EBITDA-positive "
     "month."),
    ("Q4 2028", "rocket", "320 subscribers. Marseille live, Series "
     "A conversations open."),
]
n = len(miles)
gx, gy = Inches(0.85), Inches(2.4)
cw = Inches(1.97)
# timeline
rect(s, gx + cw / 2, gy + Inches(0.46), cw * (n - 1), Inches(0.05),
     SKY)
for i, (q, ic, body) in enumerate(miles):
    x = gx + i * cw
    cx = x + cw / 2
    disc = Inches(0.94)
    oval(s, cx - disc / 2, gy, disc, disc, NAVY)
    place_icon(s, ic, cx - disc / 2 + Inches(0.25),
               gy + Inches(0.25), Inches(0.44), GOLD_HEX)
    text(s, q, x, gy + Inches(1.12), cw, Inches(0.4), size=14,
         bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    card = rrect(s, x + Inches(0.1), gy + Inches(1.55),
                 cw - Inches(0.2), Inches(2.05), WHITE, line=BORDER,
                 radius=0.1, shadow=True)
    text(s, body, x + Inches(0.27), gy + Inches(1.78),
         cw - Inches(0.54), Inches(1.7), size=10.7, color=INK,
         align=PP_ALIGN.CENTER, spacing=1.3)
# bottom strip
by = Inches(6.35)
rrect(s, Inches(0.85), by, Inches(11.65), Inches(0.62), ACCT,
      radius=0.16)
text(s, "By month 18 we will know three things: margin discipline, "
        "retention between years, and whether the brand travels "
        "between cities.",
     Inches(1.15), by, Inches(11.1), Inches(0.62), size=11,
     italic=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# =====================================================================
# SLIDE 17 — Why we will win
# =====================================================================
s = new_slide(dark=True)
rect(s, 0, 0, SW, Inches(0.14), GOLD)
text(s, "CLOSING", Inches(0.85), Inches(0.7), Inches(6),
     Inches(0.35), size=12, bold=True, color=GOLD)
text(s, "Why we will win", Inches(0.85), Inches(1.0), Inches(11),
     Inches(0.9), size=34, bold=True, color=WHITE)
rect(s, Inches(0.88), Inches(1.85), Inches(0.7), Inches(0.06), GOLD)
reasons = [
    ("users", "We are the customer",
     "Between the seven of us we lived this problem. Every feature "
     "started as something one of us needed on arrival."),
    ("lock", "The moat compounds",
     "Verified inventory, an in-house guarantor and a service layer "
     "are easy to describe and hard to assemble."),
    ("compass", "We know when to stop",
     "Below 120 subscribers and above EUR 280 CAC at month 18, we "
     "wind down. The product earns its growth or it does not."),
]
cw = Inches(3.85)
gx, gy = Inches(0.85), Inches(2.4)
for i, (ic, t, b) in enumerate(reasons):
    x = gx + i * (cw + Inches(0.2))
    rrect(s, x, gy, cw, Inches(3.5), NAVY2, radius=0.07)
    rect(s, x, gy, cw, Inches(0.1), GOLD)
    icon_chip(s, ic, x + Inches(0.35), gy + Inches(0.35),
              chip=Inches(0.9), fill=NAVY,
              icon_hex=GOLD_HEX)
    text(s, t, x + Inches(0.35), gy + Inches(1.45),
         cw - Inches(0.7), Inches(0.5), size=17, bold=True,
         color=WHITE)
    text(s, b, x + Inches(0.35), gy + Inches(2.0),
         cw - Inches(0.7), Inches(1.4), size=12, color=SKY,
         spacing=1.35)
text(s, "EURONEST", Inches(0.85), Inches(6.55), Inches(4),
     Inches(0.35), size=10, bold=True, color=SKY)
text(s, "Student Housing & Lifestyle Subscription  ·  April 2026",
     SW - Inches(6.0), Inches(6.55), Inches(5.15), Inches(0.35),
     size=10, color=SKY, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 18 — Closing
# =====================================================================
s = new_slide(dark=True)
place_photo(s, "closing", 0, 0, SW, SH, overlay=True,
            label="Closing photo (full-bleed)",
            spec="Suggested: a wide, warm shot of a city or a "
                 "student space. Landscape, min 2400x1350px.")
# dark scrim for text legibility
scrim = rect(s, 0, 0, SW, SH, NAVY)
_set_alpha(scrim, 32)
rect(s, 0, SH - Inches(0.14), SW, Inches(0.14), GOLD)
text(s, "Thank you.", Inches(0.9), Inches(2.35), Inches(11),
     Inches(1.5), size=68, bold=True, color=WHITE)
rect(s, Inches(0.95), Inches(3.85), Inches(0.85), Inches(0.07), GOLD)
text(s, "We would like to hear your questions.", Inches(0.95),
     Inches(4.1), Inches(11), Inches(0.6), size=22, italic=True,
     color=WHITE)
# contact chips
for i, (ic, t) in enumerate([("chat", "founders@euronest.eu"),
                              ("pin", "ISTEC, Paris")]):
    x = Inches(0.95) + i * Inches(3.7)
    rrect(s, x, Inches(5.15), Inches(3.4), Inches(0.66), NAVY2,
          radius=0.3)
    place_icon(s, ic, x + Inches(0.2), Inches(5.28), Inches(0.4),
               GOLD_HEX)
    text(s, t, x + Inches(0.72), Inches(5.15), Inches(2.6),
         Inches(0.66), size=13, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
out = os.path.join(HERE, "Euronest_Pitch_Deck.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
